'''
Created on FEB 2, 2020

@author: DashA2
'''
import pandas as pd
import datetime as dt
from datetime import datetime
import time
import os
from os import listdir
from os.path import isfile, join
from openpyxl.formatting.rule import ColorScaleRule
import numpy as np
from pandas import ExcelWriter
import re
from collections import Counter

import pickle4reducer
import multiprocessing as multiprocessing
ctx = multiprocessing.get_context()
ctx.reducer = pickle4reducer.Pickle4Reducer()
from functools import partial
import json
import glob
from pptx.dml.color import RGBColor
import sys


#ppt
from pptx import Presentation, opc
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches,Pt

from random import randint
from calendar import month
PP_ALIG = PP_PARAGRAPH_ALIGNMENT()

# This function writes heatmaps to excel .csv files and creates different tabs for each VNF
def write_vnf_to_excel(value_ds, excel_writer, vnf_name_init, rule=None, rule2=None):
    vnf_name = vnf_name_init.split()[0]
    if vnf_name == 'Discovered':
        vnf_name = vnf_name_init.split()[1]
    if len(vnf_name) > 30:
        vnf_name = vnf_name[:29]
    # make sure the VNF name has no strange characters, if you see an error, could be this
    #vnf_name = vnf_name.replace("/", "-").replace("\\", "-").replace(":", "-").replace("_", "-")
    #vnf_name = vnf_name.upper()
    workbook = excel_writer.book
    worksheets = workbook.worksheets
    if vnf_name in worksheets:
        for i in range(1, 1000):
            if vnf_name + str(i) in worksheets:
                vnf_name = vnf_name + str(i)
                break
    try:
        value_ds.to_excel(excel_writer, sheet_name=vnf_name)           
        if not rule == None:
            ws = workbook[vnf_name]
            ws.conditional_formatting.add('B2:BA5000', rule)
            ws.column_dimensions["A"].width = 40
    except Exception:
        print('Error Writing to file!!')
    #save_close_writer(excel_writer)
    return


# This function writes the BU Sumamry data into one csv and separates by tabs
def write_summary_to_excel(value_ds, excel_writer, counter, rule=None, rule2=None):
    workbook = excel_writer.book
    BH_cont_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                              mid_type='num', mid_value=10, mid_color='ffff00',
                              end_type='num', end_value=10, end_color='ffff00')
    BH_latency_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                 mid_type='num', mid_value=15, mid_color='ffff00',
                                 end_type='num', end_value=15, end_color='ffff00')
    BH_Stop_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                  mid_type='num', mid_value=5, mid_color='ffff00',
                                  end_type='num', end_value=5, end_color='ffff00')
    try:
        value_ds.to_excel(excel_writer, sheet_name=counter)
        if rule == BH_cont_rule or rule == BH_latency_rule or rule == BH_Stop_rule:
            ws = workbook[counter]
            # for the contention/ready sumamry page apply to colourscale rules for different cells 
            ws.conditional_formatting.add('B2:D500', rule)
            ws.conditional_formatting.add('E2:G500', rule2)
            ws.column_dimensions["A"].width = 20
            ws.column_dimensions["B"].width = 20
            ws.column_dimensions["C"].width = 20
            ws.column_dimensions["D"].width = 20
            ws.column_dimensions["E"].width = 20
            ws.column_dimensions["F"].width = 20
            ws.column_dimensions["G"].width = 20
            ws.column_dimensions["H"].width = 20
            ws.column_dimensions["I"].width = 20
            ws.row_dimensions[1].height = 60
            ws.add_format({'align': 'center'})
        if not rule == None:
            # counter is the KPI you put in, it would just create an extra tab in excel
            ws = workbook[counter]
            ws.column_dimensions["A"].width = 20
            ws.column_dimensions["B"].width = 20
            ws.column_dimensions["C"].width = 20
            ws.column_dimensions["D"].width = 20
            ws.column_dimensions["E"].width = 20
            ws.column_dimensions["F"].width = 20
            ws.column_dimensions["G"].width = 20
            ws.column_dimensions["H"].width = 20
            ws.column_dimensions["I"].width = 20
            ws.row_dimensions[1].height = 60
            ws.conditional_formatting.add('C2:I500', rule)
            ws.write('A1', 'VNF')
    except Exception:
        print('Exception while applying some styles to worksheet !!')
    return

# This function sves and closes to csv, make sure there is data in csv before saving and closing
# else the file created will be corrupt.
def save_close_writer(excel_writer):
    excel_writer.save()
    excel_writer.close()

# This fucntion get all files in the folder
def get_excel_files(path):
    return [f for f in listdir(path) if isfile(join(path, f))]

def count_values_in_range(series, range_min, range_max):
    # "between" returns a boolean Series equivalent to left <= series <= right.
    # NA values will be treated as False.
    return series.between(left=range_min, right=range_max).sum()

dirpath = os.getcwd()

if not os.path.exists(os.path.join(dirpath,'OPCO_OUTPUT')):
    os.mkdir(os.path.join(dirpath,'OPCO_OUTPUT'))
OutputFolder = os.path.join(dirpath,'OPCO_OUTPUT')
def executeVNF(vnf,maxavg,ds_vm,min_date,max_date,site):  
    bh_summary = pd.DataFrame(
        columns=["Busiest Hour (of the month)", "Busiest Hour CPU Demand (%)[Avg of all VMs] ",
                 "Busiest Hour CPU Demand (%)[Busiest VM]", "Most Common Busy Hour [MCBH] in month",
                 "MCBH CPU Demand (%) in the month[Avg of all VM]"])
    bh_summary_RAM = pd.DataFrame(
        columns=["Busiest Hour (of the month)", "Busiest Hour Memory Demand   (%)[Avg of all VMs] ",
                 "Busiest Hour CRAM Demand   (%)[Busiest VM]", "Most Common Busy Hour [MCBH] in month",
                 "MCBH Memory Demand  (%) in the month[Avg of all VM]"])
    bh_summary_cont = pd.DataFrame(columns=["Busiest Hour CPU Contention   (%)[Avg of all VMs] ",
                                            "Busiest Hour CPU Contention   (%)[Busiest VM]",
                                            "MCBH CPU Contention  (%) in the month[Avg of all VMs]"
        , "Busiest Hour CPU Ready   (%)[Avg of all VMs] ", "Busiest Hour CPU Ready   (%)[Busiest VM] ",
                                            "MCBH CPU Ready  (%) in the month[Avg of all VMs]"])

    bh_summary_latency = pd.DataFrame(
        columns=[" Storage Latency Monthly Peak Avg", "Busiest Hour Storage Latency   (ms)[Avg of all VMs] ",
                 "Busiest Hour Storage Latency   (ms)[Busiest VM]",
                 " NW Packet Drop Monthly Peak Avg", "Busiest Hour NW Packet Drop (%)[Avg of all VMs] ",
                 "Busiest Hour NW Packet Drop   (%)[Busiest VM] "])
    bh_summary_stop = pd.DataFrame(columns=["Busiest Hour CPU Co-Stop   (%)[Avg of all VMs] ",
                                            "Busiest Hour CPU Co-Stop   (%)[Busiest VM]",
                                            "MCBH CPU Co-Stop  (%) in the month[Avg of all VMs]"])
    bh_summary_at_risk = pd.DataFrame(
        columns=["Total number of VMs", "CPU > 80% ", "Memory > 80%", "Contention > 5%","Co-Stop > 5%", "CPU Ready > 5%",
                 "Storage Latency > 15ms", "NW Packet Drop > 1%"])
    print("Running VNF  "+vnf +" of "+maxavg + " Sheet")
    critical_CPU_demand = 0
    critical_CPU_ready = 0
    critical_CPU_contention = 0
    critical_NW_packet_drop = 0
    critical_latency = 0
    critical_ram = 0
    
    vnf_ds = ds_vm[ds_vm['VNF'] == vnf]
    vm_set = set(vnf_ds['VM'])
    
    cpu_demand_bh = pd.DataFrame()
    cpu_ready_bh = pd.DataFrame()
    latency_bh = pd.DataFrame()
    cpu_contention_bh = pd.DataFrame()
    cpu_stop_bh  = pd.DataFrame()
    nw_pckt_drop = pd.DataFrame()
    ram_bh = pd.DataFrame()
    

    cpu_demand_bh_c = pd.DataFrame()
    cpu_ready_bh_c = pd.DataFrame()
    latency_bh_c = pd.DataFrame()
    cpu_contention_bh_c = pd.DataFrame()
    nw_pckt_drop_c = pd.DataFrame()
    ram_bh_c = pd.DataFrame()

    date = min_date
    temp_maxA = pd.DataFrame()        
    #get the mid of the month
    mid_date = max_date - dt.timedelta(days=15)
    # Now in this while loop we go through the days of the month in the file 
    
    BH_Month_VM= pd.DataFrame()
    critical_CPU_demand_file = 0
    critical_CPU_ready_file = 0
    critical_CPU_contention_file = 0
    critical_NW_packet_drop_file = 0
    critical_latency_file = 0
    critical_ram_file = 0
    
    while date <= max_date:
        # each time in the loop, go +24 hours
        next_date = date + dt.timedelta(hours=24)
        day_hour_set = set(vnf_ds[(vnf_ds['timestamp'] >= date) & (vnf_ds['timestamp'] < next_date)]['timestamp'])
        # Here we find the Busy hour by looking into every hour of every day
        currentMonth = mid_date.strftime("%m") #4
        dateMonth = date.strftime("%m")#5
#         if not currentMonth == dateMonth:
#             print("Removing data of the month:"+ dateMonth +"from the data")
#             day_hour_set.clear()
        BH_day_Hr_df = pd.DataFrame([])
        temp_max_All = pd.DataFrame()
        if len(day_hour_set) > 0:
            for tstamp in day_hour_set:
                VM_df_temp = vnf_ds[vnf_ds['timestamp'] == tstamp]
                if WeightedVNF == True :
                    avg_cpu_value = (((VM_df_temp['cpu|Demand (%)'])*VM_df_temp['cpu|vCPU Cores']).sum()/VM_df_temp['cpu|vCPU Cores'].sum())
                else:
                    avg_cpu_value = np.nanmean((VM_df_temp['cpu|Demand (%)']).dropna())
                avg_con_value = np.nanmean((VM_df_temp['cpu|Contention (%)']).dropna())
                avg_ready_value = np.nanmean((VM_df_temp['cpu|Ready (%)']).dropna())
                avg_mem_value = np.nanmean((VM_df_temp['mem|Usage %']).dropna())
                avg_stop_value = np.nanmean((VM_df_temp['cpu|Co-Stop (%)']).dropna())
                data = {'timestamp':[tstamp],'cpu|Demand (%)':[avg_cpu_value],'mem|Usage %':[avg_mem_value],'cpu|Contention (%)':[avg_con_value],
                        'cpu|Ready (%)':[avg_ready_value],'cpu|Co-Stop (%)':[avg_stop_value]}  
                VM_df = pd.DataFrame(data)
                BH_day_Hr_df = BH_day_Hr_df.append(VM_df)#avg of Vms values for every timestamp

                temp_maxA = vnf_ds.filter(['timestamp','VNF','VM','cpu|Demand (%)','cpu|Contention (%)', 
                                          'cpu|Ready (%)', 'net:Aggregated Packet Drop Rate (%)',
                                           "mem|Usage %", 'storage|Write Latency (ms)',
                                           'storage|Read Latency (ms)','cpu|Co-Stop (%)','cpu|vCPU Cores']) #only required columns dataframe
                temp_maxA = temp_maxA.loc[temp_maxA['timestamp'] == tstamp]# for each timestamp for all vms
                temp_max_All = pd.concat([temp_max_All, temp_maxA])

            for vm in vm_set:
                temp_max_VM = temp_max_All
                day = tstamp.strftime("%d-%m")
                temp_max_VM = temp_max_VM.loc[temp_max_VM['VM'] == vm]
                
                #Highest value and date -CPU
                if not temp_max_VM.empty:
                    temp_max_VM = temp_max_VM.sort_values('cpu|Demand (%)', ascending=False)
                    highest_value_date = temp_max_VM.iloc[0,0]
                    temp_max_VM_ram = temp_max_VM.sort_values('mem|Usage %', ascending=False)
                    highest_value_DateRAM = temp_max_VM_ram.iloc[0,0]
                    
                    temp_max_VM_cont = temp_max_VM.sort_values('cpu|Contention (%)', ascending=False)
                    highest_value_Datecont = temp_max_VM_cont.iloc[0,0]
                    #cpu|Co-Stop (%)
                    temp_max_VM_Stop = temp_max_VM.sort_values('cpu|Co-Stop (%)', ascending=False)
                    highest_value_Date_Stop = temp_max_VM_Stop.iloc[0,0]
                    
                    temp_max_VM_ready = temp_max_VM.sort_values('cpu|Ready (%)', ascending=False)
                    highest_value_Dateready = temp_max_VM_ready.iloc[0,0]
                    
                    temp_max_VM_packet = temp_max_VM.sort_values('net:Aggregated Packet Drop Rate (%)', ascending=False)
                    highest_value_date_packet = temp_max_VM_packet.iloc[0,0]
                    
                    temp_ds = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_date)]
                    temp_ds_memory = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_DateRAM)]
                    temp_ds_cont = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_Datecont)]
                    temp_ds_ready = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_Dateready)]
                    temp_ds_packet = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_date_packet)]
                    #cpu|Co-Stop (%)
                    temp_ds_Stop = temp_max_VM[(temp_max_VM['timestamp'] == highest_value_Date_Stop)]
                    latency_bh['VNF'] = vnf
                    if not temp_ds['storage|Write Latency (ms)'].isnull().values.any():
                        
                        latency_bh.at[vm, day] = np.round((temp_ds['storage|Write Latency (ms)'].iloc[0] +
                                                    temp_ds['storage|Read Latency (ms)'].iloc[0]), 2)
                        if latency_bh.at[vm, day] > 14.999:
                            critical_latency = critical_latency + 1
                            critical_latency_file = critical_latency_file + 1
                    cpu_demand_bh['VNF'] = vnf
                    if not temp_ds['cpu|Demand (%)'].isnull().values.any():
                        
                        cpu_demand_bh.at[vm, day] = np.round(temp_ds['cpu|Demand (%)'].iloc[0], 0)
                        
                        if cpu_demand_bh.at[vm, day] > 79.999:
                            critical_CPU_demand = critical_CPU_demand + 1
                            critical_CPU_demand_file = critical_CPU_demand_file + 1
                    cpu_contention_bh['VNF'] = vnf  
                    if not temp_ds_cont['cpu|Contention (%)'].isnull().values.any(): 
                        
                        cpu_contention_bh.at[vm, day ] = np.round(temp_ds_cont['cpu|Contention (%)'].iloc[0], 2)
                        if cpu_contention_bh.at[vm, day] > 4.999:
                            critical_CPU_contention = critical_CPU_contention + 1
                            critical_CPU_contention_file = critical_CPU_contention_file + 1
                    cpu_ready_bh['VNF'] = vnf
                    if not temp_ds_ready['cpu|Ready (%)'].isnull().values.any():
                        
                        cpu_ready_bh.at[vm, day] = np.round(temp_ds_ready['cpu|Ready (%)'].iloc[0], 2)
                        if cpu_ready_bh.at[vm, day] > 4.999:
                            critical_CPU_ready = critical_CPU_ready + 1
                            critical_CPU_ready_file = critical_CPU_ready_file + 1
                    nw_pckt_drop['VNF'] = vnf
                    if not temp_ds_packet['net:Aggregated Packet Drop Rate (%)'].isnull().values.any():
                        
                        nw_pckt_drop.at[vm, day] = np.round(temp_ds_packet['net:Aggregated Packet Drop Rate (%)'].iloc[0], 2)
                        if nw_pckt_drop.at[vm, day] > 0.999:
                            critical_NW_packet_drop = critical_NW_packet_drop + 1
                            critical_NW_packet_drop_file = critical_NW_packet_drop_file + 1
                    ram_bh['VNF'] = vnf
                    if not temp_ds_memory["mem|Usage %"].isnull().values.any():
                        
                        ram_bh.at[vm, day] = np.round(temp_ds_memory["mem|Usage %"].iloc[0], 2)
                        if ram_bh.at[vm, day] > 79.999:
                            critical_ram = critical_ram + 1
                            critical_ram_file = critical_ram_file + 1
                    cpu_stop_bh['VNF'] = vnf
                    if not temp_ds_Stop['cpu|Co-Stop (%)'].isnull().values.any():
                        
                        cpu_stop_bh.at[vm, day] = np.round(temp_ds_Stop['cpu|Co-Stop (%)'].iloc[0], 2)
                        
                    date = next_date
            # once out of the loop, we write those criticalities to file for each vnf
            #once out of the loop have the data from
            #sort to get the peak value of the current VM
    
            if critical_CPU_demand > 0:
                cpu_demand_bh_c = cpu_demand_bh
            if critical_CPU_ready > 0:
                cpu_ready_bh_c = cpu_ready_bh
            if critical_CPU_contention > 0:
                cpu_contention_bh_c = cpu_contention_bh
            if critical_NW_packet_drop > 0:
                nw_pckt_drop_c = nw_pckt_drop
            if critical_latency > 0:
                latency_bh_c = latency_bh
            if critical_ram > 0:
                ram_bh_c = ram_bh
            try:
                vnf_n = vnf.split()[0]
            except:
                vnf_n = vnf
        date = next_date
        BH_Month_VM = BH_Month_VM.append(BH_day_Hr_df)
        
    if not BH_Month_VM.empty:
        BH_Month_VM = BH_Month_VM[['timestamp','cpu|Demand (%)','mem|Usage %','cpu|Contention (%)','cpu|Ready (%)','cpu|Co-Stop (%)']]
    BH_Month_VM.reset_index(inplace = True, drop = True) 
    critical_heat_map_dict = {'cpu_demand_bh_c':cpu_demand_bh_c,'cpu_ready_bh_c':cpu_ready_bh_c,'cpu_contention_bh_c':cpu_contention_bh_c,'nw_pckt_drop_c':nw_pckt_drop_c,'latency_bh_c':latency_bh_c,'ram_bh_c':ram_bh_c}

    ds_vm_vnf = ds_vm[ds_vm['VNF'] == vnf]
    
    ds_vm_cpu = ds_vm_vnf.copy()
    
    BH_Month_VM = BH_Month_VM.sort_values('cpu|Demand (%)', ascending=False)
    
    #Column B
    BH_ofThemonth = pd.to_datetime(BH_Month_VM.iloc[0,0]).strftime("%H:00 -> %H:59 (%m-%d)")#get the BH date
    #Column C
    Busiest_Hour_CPU_Demand = np.round(BH_Month_VM.iloc[0,1],2)
    #Column D
    Busy_hour_date_cpu = BH_Month_VM.iloc[0,0]
    ds_vm_df  = ds_vm_cpu[ds_vm_cpu['timestamp'] == Busy_hour_date_cpu]
    
    ds_vm_df = ds_vm_df.sort_values('cpu|Demand (%)', ascending=False)
    #bst_cpu_peak = np.round(ds_vm_df.iloc[0,10],2)#check sheet column number
    #ds_vm_df=ds_vm_df.reset_index(drop=True)
    #ds_vm_df['index'] = np.arange(len(ds_vm_df))
    ds_vm_df.reset_index(inplace = True, drop = True)
    s = ds_vm_df['cpu|Demand (%)']
    bst_cpu_peak = np.round(s.iloc[0],2)
    #bst_cpu_peak = np.round(ds_vm_df.at[0,'cpu|Demand (%)'],2)
    #Column E:
    Mcbh_dic = MCBH_Cal(BH_Month_VM,'cpu|Demand (%)')
    for key, val in Mcbh_dic.items():
        if key=='MCBHHR':
            MCBH_Hr_Cpu = val
        elif key == 'MCBHAvg':
            MCBH_avg_cpu = val
    #Memory/RAM:~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    ds_vm_mem = ds_vm_vnf.copy()
    #Column B:
    BH_Month_VM = BH_Month_VM.sort_values('mem|Usage %', ascending=False)
    BH_ofThemonth_ram = BH_Month_VM.iloc[0,0].strftime("%H:00 -> %H:59 (%m-%d)")#get the BH date
    #Column C:
    Busiest_Hour_RAM_Demand = np.round(BH_Month_VM.iloc[0,2],2)
    #Column D
    Busy_hour_date_mem = BH_Month_VM.iloc[0,0]
    ds_vm_df  = ds_vm_mem[ds_vm_cpu['timestamp'] == Busy_hour_date_mem]
    ds_vm_df = ds_vm_df.sort_values('mem|Usage %', ascending=False)
    #bst_ram_peak = np.round(ds_vm_df.iloc[0,14],2)#check sheet column number
    #ds_vm_df=ds_vm_df.reset_index(drop=True)
    #ds_vm_df['index'] = np.arange(len(ds_vm_df))
    ds_vm_df.reset_index(inplace = True, drop = True)
    bst_ram_peak = np.round(ds_vm_df.at[0,'mem|Usage %'],2)#check sheet column number
    #Column E & F
    Mcbh_dic = MCBH_Cal(BH_Month_VM,'mem|Usage %')
    for key, val in Mcbh_dic.items():
        if key=='MCBHHR':
            MCBH_Hr_Memory = val
        elif key == 'MCBHAvg':
            MCBH_avg_mem = val

    #CPU contention: ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    ds_vm_count = ds_vm_vnf.copy() 
    BH_Month_VM = BH_Month_VM.sort_values('cpu|Contention (%)', ascending=False)
    #Column B: C
    Busiest_Hour_cont_Demand = np.round(BH_Month_VM.iloc[0,3],2)
    #Column C : D
    Busy_hour_date_cont = BH_Month_VM.iloc[0,0]
    ds_vm_df  = ds_vm_count[ds_vm_cpu['timestamp'] == Busy_hour_date_cont]
    ds_vm_df = ds_vm_df.sort_values('cpu|Contention (%)', ascending=False)
    #bst_cont_peak = np.round(ds_vm_df.iloc[0,11],2)#check sheet column number
    #ds_vm_df=ds_vm_df.reset_index(drop=True)
    #ds_vm_df['index'] = np.arange(len(ds_vm_df))
    ds_vm_df.reset_index( inplace = True, drop = True)
    bst_cont_peak = np.round(ds_vm_df.at[0,'cpu|Contention (%)'],2)
    #Column D: D & G
    Mcbh_dic = MCBH_Cal(BH_Month_VM,'cpu|Contention (%)')
    for key, val in Mcbh_dic.items():
        if key=='MCBHHR':
            MCBH_Hr_Memory = val
        elif key == 'MCBHAvg':
            Peak_cont = val
    #CPU READY ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    ds_vm_ready = ds_vm_vnf.copy()
    BH_Month_VM = BH_Month_VM.sort_values('cpu|Ready (%)', ascending=False)
    #Column E:C
    Busiest_Hour_cpu_ready = np.round(BH_Month_VM.iloc[0,4],2)
    #Column F:D
    Busy_hour_date = BH_Month_VM.iloc[0,0]#get BH values
    ds_vm_df  = ds_vm_ready[ds_vm_ready['timestamp'] == Busy_hour_date]
    ds_vm_df = ds_vm_df.sort_values('cpu|Ready (%)', ascending=False)
    #bst_ready_peak = np.round(ds_vm_df.iloc[0,12],2)#check sheet column number
    #ds_vm_df=ds_vm_df.reset_index(drop=True)
    #ds_vm_df['index'] = np.arange(len(ds_vm_df))
    ds_vm_df.reset_index(inplace = True, drop = True)
    bst_ready_peak = np.round(ds_vm_df.at[0,'cpu|Ready (%)'],2)
    #Column G: F
    Mcbh_dic = MCBH_Cal(BH_Month_VM,'cpu|Ready (%)')
    for key, val in Mcbh_dic.items():
        if key=='MCBHHR':
            MCBH_Hr_Memory = val
        elif key == 'MCBHAvg':
            Peak_ready = val
    #CO -STOP: ~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    ds_vm_stop = ds_vm_vnf.copy()
    BH_Month_VM = BH_Month_VM.sort_values('cpu|Co-Stop (%)', ascending=False)
    #Column B:
    bst_cont_peak_stop = np.round(BH_Month_VM.iloc[0,5],2)
    #Column C:
    Busy_hour_date = BH_Month_VM.iloc[0,0]#get BH values
    ds_vm_df  = ds_vm_ready[ds_vm_stop['timestamp'] == Busy_hour_date]
    ds_vm_df = ds_vm_df.sort_values('cpu|Co-Stop (%)', ascending=False)
    #Busiest_Hour_Stop_Demand = np.round(ds_vm_df.iloc[0,13],2)#check sheet column number
    #ds_vm_df=ds_vm_df.reset_index(drop=True)
    #ds_vm_df['index'] = np.arange(len(ds_vm_df))
    ds_vm_df.reset_index(inplace = True, drop = True)
    Busiest_Hour_Stop_Demand = np.round(ds_vm_df.at[0,'cpu|Co-Stop (%)'],2)#check sheet column number
    #Column D:

    Mcbh_dic = MCBH_Cal(BH_Month_VM,'cpu|Co-Stop (%)')
    for key, val in Mcbh_dic.items():
        if key=='MCBHHR':
            MCBH_Hr_Memory = val
        elif key == 'MCBHAvg':
            Peak_stop = val
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    #latency sheet
    Max_Month_array = ds_vm[ds_vm['VNF'] == vnf]
    latency_month_peak_write = Max_Month_array['storage|Write Latency (ms)'].mean()
    latency_month_peak_read = Max_Month_array['storage|Read Latency (ms)'].mean()
    latency_month_peak = np.round((latency_month_peak_write + latency_month_peak_read) / 2, 2)
    
    packet_drop_month_peak = np.round(Max_Month_array['net:Aggregated Packet Drop Rate (%)'].mean(), 2)
    
    RAM_Array = Max_Month_array[Max_Month_array['timestamp'] == Busy_hour_date_cpu]# busiest hour by CPU demand
    Busiest_Hour_latency_read = RAM_Array['storage|Read Latency (ms)'].mean()
    Busiest_Hour_latency_write = RAM_Array['storage|Write Latency (ms)'].mean()
    Busiest_Hour_latency = np.round((Busiest_Hour_latency_read + Busiest_Hour_latency_write) / 2, 2)
    
    bst_latency_read_peak = RAM_Array['storage|Read Latency (ms)'].max()
    bst_latency_write_peak = RAM_Array['storage|Write Latency (ms)'].max()
    bst_latency_peak = np.round((bst_latency_read_peak + bst_latency_write_peak) / 2, 2)
    
    Busiest_Hour_packet_drop = np.round(RAM_Array['net:Aggregated Packet Drop Rate (%)'].mean(), 2)
    
    bst_packet_drop_peak = np.round(RAM_Array['net:Aggregated Packet Drop Rate (%)'].max(), 2)
    
    # create new columns in the dataframe 
    cpu_demand_bh['Mean'] = np.round(cpu_demand_bh.mean(axis=1), 2)
    cpu_demand_bh['Max'] = np.round(cpu_demand_bh.max(axis=1), 2)
    #cpu_demand_bh= getMaxTimestamp(cpu_demand_bh,vnf_ds,'cpu|Demand (%)')

    # here you can define the boundary and count
    cpu_demand_bh["# Days Above 80%"] = cpu_demand_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 80, 200), axis=1)

    cpu_contention_bh['Mean'] = np.round(cpu_contention_bh.mean(axis=1), 2)
    cpu_contention_bh['Max'] = np.round(cpu_contention_bh.max(axis=1), 2)
    cpu_contention_bh["# Days Above 5%"] = cpu_contention_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 5, 100), axis=1)
    
    cpu_stop_bh['Mean'] = np.round(cpu_stop_bh.mean(axis=1), 2)
    cpu_stop_bh['Max'] = np.round(cpu_stop_bh.max(axis=1), 2)
    cpu_stop_bh["# Days Above 5%"] = cpu_stop_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 5, 100), axis=1)

    cpu_ready_bh['Mean'] = np.round(cpu_ready_bh.mean(axis=1), 2)
    cpu_ready_bh['Max'] = np.round(cpu_ready_bh.max(axis=1), 2)
    cpu_ready_bh["# Days Above 5%"] = cpu_ready_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 5, 100), axis=1)

    ram_bh['Mean'] = np.round(ram_bh.mean(axis=1), 2)
    ram_bh['Max'] = np.round(ram_bh.max(axis=1), 2)
    ram_bh["# Days Above 80%"] = ram_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 80, 100), axis=1)

    nw_pckt_drop['Mean'] = np.round(nw_pckt_drop.mean(axis=1), 2)
    nw_pckt_drop['Max'] = np.round(nw_pckt_drop.max(axis=1), 2)
    nw_pckt_drop["# Days Above 1%"] = nw_pckt_drop.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 1, 100), axis=1)

    latency_bh['Mean'] = np.round(latency_bh.mean(axis=1), 2)
    latency_bh['Max'] = np.round(latency_bh.max(axis=1), 2)
    latency_bh["# Days Above 15ms"] = latency_bh.drop(["Mean", "Max","VNF"], axis=1).apply(
        func=lambda row: count_values_in_range(row, 15, 200), axis=1)    
    heat_map_dict = {'latency_bh':latency_bh,'cpu_contention_bh':cpu_contention_bh,'nw_pckt_drop':nw_pckt_drop,'cpu_demand_bh':cpu_demand_bh,'cpu_ready_bh':cpu_ready_bh,'ram_bh':ram_bh,'cpu_stop_bh':cpu_stop_bh}

    total_vm_df = ds_vm[ds_vm['VNF'] == vnf]
    total_vm_set = set(total_vm_df['VM'])
    #CPU > 80%
    above_threshold_count_CPU_demand_list = cpu_demand_bh['# Days Above 80%'].tolist()
    above_threshold_count_CPU_demand = len([i for i in above_threshold_count_CPU_demand_list if i > 0]) 
    #Memory > 80%
    above_threshold_count_RAM_list = ram_bh["# Days Above 80%"].tolist()
    above_threshold_count_RAM = len([i for i in above_threshold_count_RAM_list if i > 0])
    #Contention > 5%
    above_threshold_count_CPU_contention_list = cpu_contention_bh["# Days Above 5%"].tolist()
    above_threshold_count_CPU_contention = len([i for i in above_threshold_count_CPU_contention_list if i > 0])
    #Contention > 5%
    above_threshold_count_CPU_stop_list = cpu_stop_bh["# Days Above 5%"].tolist()
    above_threshold_count_CPU_stop = len([i for i in above_threshold_count_CPU_stop_list if i > 0])
    
    #CPU Ready > 5%
    above_threshold_count_CPU_ready_list = cpu_ready_bh["# Days Above 5%"].tolist()
    above_threshold_count_CPU_ready = len([i for i in above_threshold_count_CPU_ready_list if i > 0])
    #Storage Latency > 15ms
    above_threshold_count_latencylist = latency_bh["# Days Above 15ms"].tolist()
    above_threshold_count_latency = len([i for i in above_threshold_count_latencylist if i > 0])
    #NW Packet Drop > 1%
    above_threshold_count_packet_drop_list = nw_pckt_drop["# Days Above 1%"].tolist()
    above_threshold_count_packet_drop = len([i for i in above_threshold_count_packet_drop_list if i > 0])
    
    bh_summary.at[vnf_n] = (BH_ofThemonth, Busiest_Hour_CPU_Demand, bst_cpu_peak, MCBH_Hr_Cpu, MCBH_avg_cpu)
    bh_summary_RAM.at[vnf_n] = (BH_ofThemonth_ram, Busiest_Hour_RAM_Demand, bst_ram_peak, MCBH_Hr_Memory, MCBH_avg_mem)
    bh_summary_cont.at[vnf_n] = (
        Busiest_Hour_cont_Demand, bst_cont_peak, Peak_cont, Busiest_Hour_cpu_ready, bst_ready_peak, Peak_ready)
    bh_summary_latency.at[vnf_n] = (
        latency_month_peak, Busiest_Hour_latency, bst_latency_peak, packet_drop_month_peak,
        Busiest_Hour_packet_drop,
        bst_packet_drop_peak)
    bh_summary_stop.at[vnf_n] = (bst_cont_peak_stop,Busiest_Hour_Stop_Demand,Peak_stop)

    bh_summary_at_risk.at[vnf_n] = (len(total_vm_set), above_threshold_count_CPU_demand, above_threshold_count_RAM, 
                                    above_threshold_count_CPU_contention,above_threshold_count_CPU_stop, above_threshold_count_CPU_ready 
                                        , above_threshold_count_latency, above_threshold_count_packet_drop)
    #filling Nan values fo VMs with zero: rare case
    bh_summary = bh_summary.fillna(0)
    bh_summary_RAM = bh_summary_RAM.fillna(0)
    bh_summary_cont = bh_summary_cont.fillna(0)
    bh_summary_latency = bh_summary_latency.fillna(0)
    bh_summary_stop = bh_summary_stop.fillna(0)
    bh_summary_at_risk = bh_summary_at_risk.fillna(0)
    Max_Avg_Dic = dict()
    if maxavg == 'MAX':
        max_max = np.round(cpu_demand_bh['Max'].max(), 2) #MAx of MAX column from MAX sheet
        Max_DF_Temp = pd.DataFrame({'VNF': vnf,'MAX': max_max}, index=[0])
        Max_Avg_Dic = {'MAX':Max_DF_Temp}

    if maxavg == 'AVG':
        avg_mean = np.round(np.nanmean(cpu_demand_bh['Mean']), 2) # avg of avg column from AVG sheet
        Avg_DF_Temp = pd.DataFrame({'VNF': vnf, 'AVG': avg_mean}, index=[0])
        Max_Avg_Dic = {'AVG':Avg_DF_Temp}
    
    List_Summary = {'bh_summary':bh_summary,'bh_summary_RAM':bh_summary_RAM,'bh_summary_cont':bh_summary_cont,'bh_summary_latency':bh_summary_latency,'bh_summary_stop':bh_summary_stop, 'bh_summary_at_risk':bh_summary_at_risk}
    crtical_file_list = {'critical_CPU_demand_file':critical_CPU_demand_file,'critical_latency_file':critical_latency_file,'critical_CPU_contention_file':critical_CPU_contention_file,'critical_NW_packet_drop_file':critical_NW_packet_drop_file,'critical_CPU_ready_file':critical_CPU_ready_file,'critical_ram_file':critical_ram_file}

    heatmap_data = {'heat_map_list':heat_map_dict,'critical_heat_map_list':critical_heat_map_dict}
    final_dict = [vnf,List_Summary,crtical_file_list,heatmap_data]
    return Max_Avg_Dic,final_dict
    #end of Heat map creation:
    
#fuction for calculating Most common busy hour and avg values of it.
def MCBH_Cal(BH_Month_VM, flag):
    BH_Month_VM_temp = BH_Month_VM.copy() 
    BH_Month_VM_day = pd.DataFrame()
    MCBH_Month_VM_peak = pd.DataFrame()
    set_date = set(pd.to_datetime(BH_Month_VM_temp['timestamp'].astype(str)).apply(lambda x: x.replace()).dt.strftime('%m-%d-%y'))#30 days
    #set_date = set(pd.to_datetime(BH_Month_VM_temp['timestamp']).apply(lambda x: x.replace()).dt.strftime('%m-%d-%y'))#30 days
    for date in set_date:
        date = date + ' 00:59:59'
        date = datetime.strptime(date,'%m-%d-%y %H:%M:%S')
        hr = 0
        BH_Month_VM_day_1 = pd.DataFrame()
        while(hr<24):
            BH_Month_VM_day = BH_Month_VM_temp[BH_Month_VM_temp['timestamp'] == date] #get DF for each day
            BH_Month_VM_day_1 = BH_Month_VM_day_1.append(BH_Month_VM_day)#24 value
            date = date + dt.timedelta(hours = 1)
            hr = hr + 1
        BH_Month_VM_day_1 = BH_Month_VM_day_1.sort_values(flag, ascending=False)#
        MCBH_Month_VM_peak = MCBH_Month_VM_peak.append(BH_Month_VM_day_1.head(1),ignore_index=True)#31/30 values 
    MCBH_Month_VM_peak['timestamp'] = pd.to_datetime(MCBH_Month_VM_peak['timestamp']).apply(lambda x: x.replace()).dt.strftime('%H:00 -> %H:59') 
    MCBH_Hr = Counter(MCBH_Month_VM_peak['timestamp']).most_common(1)[0][0]
    MCBH_Month_VM_peak = MCBH_Month_VM_peak[(MCBH_Month_VM_peak['timestamp'] == MCBH_Hr)]
    MCBH_avg = np.round(np.nanmean(MCBH_Month_VM_peak[flag]),2)
    Mcbh_dic = {'MCBHHR':MCBH_Hr, 'MCBHAvg':MCBH_avg}
    return Mcbh_dic
def WriteVnfData(List_of_parameters,opco_site,site,maxavg):

    bh_summary = pd.DataFrame()
    bh_summary_RAM = pd.DataFrame()
    bh_summary_cont = pd.DataFrame()
    bh_summary_latency = pd.DataFrame()
    bh_summary_at_risk = pd.DataFrame()
    bh_summary_stop = pd.DataFrame()
    
    cpu_hm= pd.DataFrame()
    cpu_con_hm= pd.DataFrame()
    cpu_pct_hm= pd.DataFrame()
    cpu_ltc_hm= pd.DataFrame()
    cpu_rdy_hm= pd.DataFrame()
    cpu_ram_hm= pd.DataFrame()
    
    ruleCPU = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                             mid_type='num', mid_value=50, mid_color='9BC2E6',
                             end_type='num', end_value=70, end_color='1F4E78')
    ruleCPU_Ready = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                                   mid_type='num', mid_value=4, mid_color='9BC2E6',
                                   end_type='num', end_value=8, end_color='1F4E78')
    rule_packet = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                                 mid_type='num', mid_value=1, mid_color='9BC2E6',
                                 end_type='num', end_value=1, end_color='1F4E78')
    rule_contention = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                                     mid_type='num', mid_value=2, mid_color='9BC2E6',
                                     end_type='num', end_value=5, end_color='1F4E78')
    rule_latency = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                                  mid_type='num', mid_value=5, mid_color='9BC2E6',
                                  end_type='num', end_value=15, end_color='1F4E78')
    ram_rule = ColorScaleRule(start_type='num', start_value=0, start_color='DDEBF7',
                              mid_type='num', mid_value=50, mid_color='9BC2E6',
                              end_type='num', end_value=70, end_color='1F4E78')
    heatmap_writer = pd.ExcelWriter(os.path.join(OutputFolder,opco_site  +'_'+site+'_'+ maxavg+ '__HeatMaps.xlsx'), engine='openpyxl')
    cpu_c_writer = pd.ExcelWriter(os.path.join(OutputFolder,opco_site  +'_'+site+ '_'+ maxavg+'__ Critical_CPU_Demand.xlsx'), engine='openpyxl')
    ram_c_writer = pd.ExcelWriter(os.path.join(OutputFolder,opco_site  +'_'+site+'_'+ maxavg+ '__Critical_RAM_Demand.xlsx'), engine='openpyxl')
    cpu_Ready_c_writer = pd.ExcelWriter(os.path.join(OutputFolder,opco_site  +'_'+site+'_'+ maxavg+ '__Critical_CPU_ready.xlsx'), engine='openpyxl')
    contention_c_writer = pd.ExcelWriter(os.path.join(OutputFolder , opco_site +'_'+site+ '_'+ maxavg+'__Critical_CPU_Contention.xlsx'),
                                         engine='openpyxl')
    latency_c_writer = pd.ExcelWriter(os.path.join(OutputFolder , opco_site +'_'+site+ '_'+ maxavg+'__Critical_Storage_Latency.xlsx'),
                                      engine='openpyxl')
    pckt_drop_c_writer = pd.ExcelWriter(os.path.join(OutputFolder , opco_site +'_'+site+'_'+ maxavg+ '__Critical_NW_Packet_Drop.xlsx'),
                                        engine='openpyxl')
    bh_summary_writer = pd.ExcelWriter(os.path.join(OutputFolder , opco_site +'_'+site+ '_'+ maxavg+'__SUMMARY.xlsx'), engine='openpyxl')

    critical_CPU_demand_file = 0
    critical_CPU_ready_file = 0
    critical_CPU_contention_file = 0
    critical_NW_packet_drop_file = 0
    critical_latency_file = 0
    critical_ram_file = 0
    #for each VNF = each_item:
    for each_item in List_of_parameters:
        vnf =each_item[0]
        list_summary_List = each_item[1]
        crtical_file = each_item[2]
        heatmap_data = each_item[3]
        for key,data in list_summary_List.items():
            if key=="bh_summary":
                bh_summary = bh_summary.append(data)
            elif key=="bh_summary_RAM":
                bh_summary_RAM = bh_summary_RAM.append(data)
            elif key=="bh_summary_cont":
                bh_summary_cont = bh_summary_cont.append(data)
            elif key=="bh_summary_latency":
                bh_summary_latency = bh_summary_latency.append(data)
            elif key =="bh_summary_stop":
                bh_summary_stop = bh_summary_stop.append(data)
            elif key=="bh_summary_at_risk":
                bh_summary_at_risk = bh_summary_at_risk.append(data)
        for key,data in crtical_file.items():
            if key=="critical_CPU_demand_file":
                critical_CPU_demand_file = critical_CPU_demand_file+ data
            elif key=="critical_CPU_ready_file":
                critical_CPU_ready_file = critical_CPU_ready_file+ data
            elif key=="critical_CPU_contention_file":
                critical_CPU_contention_file = critical_CPU_contention_file+data
            elif key=="critical_NW_packet_drop_file":
                critical_NW_packet_drop_file = critical_NW_packet_drop_file+data
            elif key=="critical_latency_file":
                critical_latency_file = critical_latency_file+data
            elif key=="critical_ram_file":
                critical_ram_file = critical_ram_file+data
        for key,data in heatmap_data.items():
            if key =="heat_map_list":
                heat_map_dict = data
                for key,heatmap in heat_map_dict.items():
                    if key == "cpu_demand_bh":
                        cpu_hm = cpu_hm.append(heatmap)
                    elif key =="cpu_contention_bh":
                        cpu_con_hm = cpu_con_hm.append(heatmap)
                    elif key =="nw_pckt_drop":
                        cpu_pct_hm = cpu_pct_hm.append(heatmap)
                    elif key =="latency_bh":
                        cpu_ltc_hm = cpu_ltc_hm.append(heatmap)
                    elif key == "cpu_ready_bh":
                        cpu_rdy_hm = cpu_rdy_hm.append(heatmap)
                    elif key == "ram_bh":
                        cpu_ram_hm = cpu_ram_hm.append(heatmap)
            elif key =="critical_heat_map_list":
                critical_heat_map_dict = data
                for key1,critical_heat_map in critical_heat_map_dict.items():
                    #if not critical_heat_map.empty():
                    if (key1 =="latency_bh_c" and not critical_heat_map.empty):
                        #if critical_latency_file > 0:
                        write_vnf_to_excel(critical_heat_map, latency_c_writer , vnf, rule_latency)
                    elif (key1 =="cpu_contention_bh_c"and not critical_heat_map.empty):
                        #if critical_CPU_contention_file > 0:
                        write_vnf_to_excel(critical_heat_map, contention_c_writer , vnf, rule_contention)
                    elif (key1 =="nw_pckt_drop_c" and not critical_heat_map.empty):
                        #if critical_NW_packet_drop_file > 0:
                        write_vnf_to_excel(critical_heat_map, pckt_drop_c_writer, vnf, rule_packet)
                    elif (key1 == "cpu_demand_bh_c" and not critical_heat_map.empty):
                        #if  critical_CPU_demand_file > 0:
                        write_vnf_to_excel(critical_heat_map, cpu_c_writer , vnf, ruleCPU)
                    elif (key1 == "cpu_ready_bh_c" and not critical_heat_map.empty):
                        #if  critical_CPU_ready_file > 0:
                        write_vnf_to_excel(critical_heat_map, cpu_Ready_c_writer , vnf, ruleCPU_Ready)
                    elif (key1 == "ram_bh_c" and not critical_heat_map.empty):
                        #if critical_ram_file > 0:
                        write_vnf_to_excel(critical_heat_map, ram_c_writer, vnf, ram_rule)

    write_vnf_to_excel(cpu_hm, heatmap_writer,  'cpu_demand_bh', ruleCPU)
    write_vnf_to_excel(cpu_con_hm, heatmap_writer,  'cpu_contention_bh', rule_contention)
    write_vnf_to_excel(cpu_pct_hm, heatmap_writer,  'nw_pckt_drop', rule_packet)
    write_vnf_to_excel(cpu_ltc_hm, heatmap_writer,  'latency_bh', rule_latency)
    write_vnf_to_excel(cpu_rdy_hm, heatmap_writer,  'cpu_ready_bh', ruleCPU_Ready)
    write_vnf_to_excel(cpu_ram_hm, heatmap_writer,  'ram_bh', ram_rule)
    save_close_writer(heatmap_writer)    

    
    if critical_CPU_demand_file > 1:
        save_close_writer(cpu_c_writer)
        print(" Critical CPU Demand values found ")
    if critical_latency_file > 1:
        save_close_writer(latency_c_writer)
        print(" Critical Storage Latency values found")
    if critical_CPU_contention_file > 1:
        save_close_writer(contention_c_writer)
        print(" Critical CPU Contention values found")
    if critical_NW_packet_drop_file > 1:
        save_close_writer(pckt_drop_c_writer)
        print(" Critical NW Packet drop values found")
    if critical_CPU_ready_file > 1:
        save_close_writer(cpu_Ready_c_writer)
        print(" Critical CPU Ready time values found")
    if critical_ram_file > 1:
        save_close_writer(ram_c_writer)
        print(" Critical Memory Demand values found")
         
    #Summary files creation:
    risk_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',#green
                           mid_type='num', mid_value=1, mid_color='ffff00',#yellow
                           end_type='num', end_value=1, end_color='ffff00')
    BH_cpu_rule = ColorScaleRule(start_type='num', start_value=0, start_color='98fb98',
                                 mid_type='num', mid_value=50, mid_color='7cfc00',
                                 end_type='num', end_value=80, end_color='228b22')
    BH_ram_rule = ColorScaleRule(start_type='num', start_value=0, start_color='98fb98',
                                 mid_type='num', mid_value=50, mid_color='7cfc00',
                                 end_type='num', end_value=80, end_color='228b22')
    BH_cont_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                  mid_type='num', mid_value=5, mid_color='ffff00',
                                  end_type='num', end_value=5, end_color='ffff00')
    BH_ready_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                   mid_type='num', mid_value=5, mid_color='ffff00',
                                   end_type='num', end_value=5, end_color='ffff00')
    BH_latency_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                     mid_type='num', mid_value=15, mid_color='ffff00',
                                     end_type='num', end_value=15, end_color='ffff00')
    BH_pdrop_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                   mid_type='num', mid_value=1, mid_color='ffff00',
                                   end_type='num', end_value=1, end_color='ffff00')
    BH_Stop_rule = ColorScaleRule(start_type='num', start_value=0, start_color='00ff7f',
                                  mid_type='num', mid_value=5, mid_color='ffff00',
                                  end_type='num', end_value=5, end_color='ffff00')
    
    try:
        write_summary_to_excel(bh_summary, bh_summary_writer, "CPU", BH_cpu_rule)
        write_summary_to_excel(bh_summary_RAM, bh_summary_writer, "Memory", BH_ram_rule)
        write_summary_to_excel(bh_summary_cont, bh_summary_writer, "Contention_&_Ready_Time", BH_cont_rule,
                               BH_ready_rule)
        write_summary_to_excel(bh_summary_latency, bh_summary_writer, "Latency_&_Packet_Drop", BH_latency_rule,
                               BH_pdrop_rule)
        write_summary_to_excel(bh_summary_stop, bh_summary_writer, "CPU_Co-Stop", BH_Stop_rule,
                               BH_ready_rule)
        write_summary_to_excel(bh_summary_at_risk, bh_summary_writer, "VMs_at_Risk", risk_rule)
    except:
        print("No Data yet")    
    save_close_writer(bh_summary_writer) 
#Executive summary:
def executiveSum(vendor_df,Max_DF,Avg_DF,opco_site,site):

    if not Max_DF.empty and not Avg_DF.empty: 
        Merge_Max_Avg_df = pd.merge(Max_DF, Avg_DF, on='VNF' , how='inner')
        merge_df = pd.merge(vendor_df, Merge_Max_Avg_df, on='VNF' , how='inner') 
        merge_df = merge_df.drop_duplicates()
        merge_df = merge_df.dropna(subset =['Vendor', 'Program'])
        merge_df['Vendor'] = merge_df['Vendor'].str.replace(r'[^\w\s]+', '')
        merge_df['Program'] = merge_df['Program'].str.replace(r'[^\w\s]+', '')
        vendor_set = set(merge_df['Vendor'].str.strip().tolist())
        vendor_set.discard('')
        vendor_list = merge_df['Vendor'].tolist()
        program_set = set(merge_df['Program'].tolist())
        program_set.discard('')
        program_list = merge_df['Program'].tolist()
        list_of_vendors_max = []
        list_of_vendors_avg = []
        Write_Dataframe_avg = pd.DataFrame({'key':['x']})
        Write_Dataframe_max = pd.DataFrame({'key':['x']})
        Write_Dataframe_avgP = pd.DataFrame({'key':['x']})
        Write_Dataframe_maxP = pd.DataFrame({'key':['x']})
        final_DF_towrite_avg = pd.DataFrame({'key':['x']})
        final_DF_towrite_max = pd.DataFrame({'key':['x']})
        if not len(vendor_list) == 0 :
            list_of_all_datas_avg = []
            list_of_all_datas_max = []
            for item in vendor_set: #each vendor names from the set.
                for item1 in vendor_list:#repeated vendors name in the list
                        if item == item1 :
                            list_of_vendors_max = merge_df.loc[merge_df['Vendor'] == item1, 'MAX'].tolist()
                            list_of_vendors_avg = merge_df.loc[merge_df['Vendor'] == item1, 'AVG'].tolist()
                avg_value = np.nanmean(list_of_vendors_avg) #for vendor 1
                max_value = np.nanmean(list_of_vendors_max)
                list_of_all_datas_avg.append(avg_value)
                list_of_all_datas_max.append(max_value)
                Write_Dataframe_avg_temp = pd.DataFrame({item:[avg_value]}) #create a dataframe
                Write_Dataframe_max_temp = pd.DataFrame({item:[max_value]})
                Write_Dataframe_avg_temp['key'] = 'x'
                Write_Dataframe_max_temp['key'] = 'x'
                  
                Write_Dataframe_avg = pd.merge(Write_Dataframe_avg, Write_Dataframe_avg_temp, how='outer', on='key')
                Write_Dataframe_max = pd.merge(Write_Dataframe_max, Write_Dataframe_max_temp, how= 'outer', on='key')
        final_DF_towrite_avg = pd.merge(final_DF_towrite_avg,Write_Dataframe_avg, how='outer',on='key')
        final_DF_towrite_max = pd.merge(final_DF_towrite_max,Write_Dataframe_max, how='outer', on='key')
          
        if not len(program_list) == 0 :
            for item in program_set: #each vendor names from the set.
                for item1 in program_list:#repeated vendors name in the list
                        if item == item1 :
                            list_of_program_max = merge_df.loc[merge_df['Program'] == item1, 'MAX'].tolist()
                            list_of_program_avg = merge_df.loc[merge_df['Program'] == item1, 'AVG'].tolist()
                avg_value = np.nanmean(list_of_program_avg) 
                max_value = np.nanmean(list_of_program_max)
                Write_Dataframe_avg_temp = pd.DataFrame({item:[avg_value]})
                Write_Dataframe_max_temp = pd.DataFrame({item:[max_value]})
                Write_Dataframe_avg_temp['key'] = 'x'
                Write_Dataframe_max_temp['key'] = 'x'
                Write_Dataframe_avgP = pd.merge(Write_Dataframe_avgP,Write_Dataframe_avg_temp, how='outer', on='key')
                Write_Dataframe_maxP = pd.merge(Write_Dataframe_maxP,Write_Dataframe_max_temp, how='outer', on='key')
        
        final_DF_towrite_avg = final_DF_towrite_avg.drop(columns=['key'])
        Write_Dataframe_avgP = Write_Dataframe_avgP.drop(columns=['key'])
        final_DF_towrite_max = final_DF_towrite_max.drop(columns=['key'])
        Write_Dataframe_maxP = Write_Dataframe_maxP.drop(columns=['key']) 
    else:
        print("******--------------------*****")
        print("Note:  Executive summary cannot be generated. Need both MAX and AVG")
    #return MAX_exe_df,AVG_exe_df
    return final_DF_towrite_avg,Write_Dataframe_avgP,final_DF_towrite_max,Write_Dataframe_maxP


#This chunksize is customized as the vnfdata can be very big and will all processes in the pool might not complete the executon of data
#on same time. Hence taking chunk size >1 will be a big disadvantage here. Taking chunksize 1 by default will also not give much advantage on execution time.
            
#PPT Defination
def ppt(series,chart_name,opco_site,name):
    cla = XL_CHART_TYPE()
    # create presentation with 1 slide ------
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = opco_site+'-'+name
    
    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = ['Total Capacity','Provisioned Capacity','Total Demand']
    chart_data.add_series('Series 1', series)
    
    # add chart to slide --------------------
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(cla.COLUMN_STACKED, x, y, cx, cy, chart_data)
    prs.save(os.path.join(OutputFolder,chart_name+'.pptx'))
#creating presentation
def colorlist(value,start_value,start_color,end_value,end_color):
    output = []
    index = 0
    if end_value==1 and start_value==0:
        end_value=10
        index = int((value/10)*100)
    else:
        value = int(np.round(value,0))
        index = value-start_value
    r1, g1, b1 = RGBColor.from_string(start_color)
    r2, g2, b2 = RGBColor.from_string(end_color)
    rdelta, gdelta, bdelta = (r2-r1)/(end_value-start_value), (g2-g1)/(end_value-start_value), (b2-b1)/(end_value-start_value)
    for step in range(0,end_value - start_value):
        r1 += rdelta
        g1 += gdelta
        b1 += bdelta
        output.append((int(r1), int(g1), int(b1)))
    color = output[index] 
    return color
       
def ColorScaleppt( value,start_value,start_color, mid_value,mid_color,end_color,end_value):
    value = int(np.round(value,0))

    if value in range(start_value,mid_value):
        return colorlist(value,start_value,start_color,mid_value,mid_color)   
    elif value in range(mid_value,end_value):
        return colorlist(value,mid_value,mid_color,end_value,end_color)
    else:
        return RGBColor.from_string(end_color)
        
def CPURAMColorScaleRule( value):
    start_value = 0 
    start_color = 'edf7ed'
    mid_value = 50
    mid_color = 'c8fac8'
    end_color = '9dfa9d'
    end_value = 80 
    return ColorScaleppt(value,start_value,start_color, mid_value,mid_color,end_color,end_value)      

def OtherColorScaleRule( sheetname,value):
    start_color = '5cfaab'
    mid_color = 'f7f781'
    end_color = 'f7f781'
    color = ''
    if sheetname in ('Contention_&_Ready_Time','CPU_Co-Stop'):
        start_value = 0
        mid_value = 5
        end_value = 5
        color = ColorScaleppt(value,start_value,start_color, mid_value,mid_color,end_color,end_value)
    elif sheetname == 'Latency' :
        start_value = 0
        mid_value = 15
        end_value = 15
        color = ColorScaleppt(value,start_value,start_color, mid_value,mid_color,end_color,end_value)
    
    elif sheetname in ('Packet_Drop','VMs_at_Risk'):
        start_value = 0
        mid_value = 1
        end_value = 1
        color = ColorScaleppt(value,start_value,start_color, mid_value,mid_color,end_color,end_value)
    return color 
def getRGBcolor( tup):
    r= int(tup[0])
    b= int(tup[1])
    g= int(tup[2])
    return RGBColor(r,b,g)
def CPUChartSlide(vnfName,vmNo,vmRisk,prs):
    cla = XL_CHART_TYPE()
    # create presentation with 1 slide ------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = 'VMs at risk - CPU Demand chart'
    #title.text_frame.fit_text(max_size=32,bold=True)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(250,0,0)
    title.text_frame.paragraphs[0].font.size=Pt(32)
    title.text_frame.paragraphs[0].font.bold=True 
    # define chart data ---------------------
    chart_data = CategoryChartData()
    
    chart_data.categories = vnfName
    #vmRisk1 = pd.Series(range(0, len(vmRisk)))#testing need to remove
    chart_data.add_series('Number of VM in VNF', vmNo)
    chart_data.add_series('CPU > 80%', vmRisk)
    
    # add chart to slide --------------------
    x, y, cx, cy = Inches(2), Inches(2), Inches(10.5), Inches(5.5)
    chart = slide.shapes.add_chart(cla.COLUMN_STACKED, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True
def DivideDf(df):
    n = 13  #chunk row size
    list_df = [df[i:i+n] for i in range(0,df.shape[0],n)]
    return list_df   

def autoSlide( dirpath,opconame,site,prs):
    sheetMap = {'CPU':'CPU Demand','Memory': 'Memory Usage', 'Contention_&_Ready_Time':'CPU Contention \ CPU Ready Time',
                'Latency':'Latency & Packet Drop','Packet_Drop':'Latency & Packet Drop','Latency_&_Packet_Drop':'Latency & Packet Drop','CPU_Co-Stop':'CPU Contention \ CPU Ready Time \ CPU Co-Stop','VMs_at_Risk':'VMs at Risk'}
    PP_ALIG = PP_PARAGRAPH_ALIGNMENT()
    #seconf slide:site slide
    if not opconame.upper() == 'AL':
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        title = slide.shapes.title
        title.text = site
        subtitle=slide.placeholders[1]
        subtitle.text = ' '
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
        title.text_frame.paragraphs[0].font.size=Pt(44)
        
    vnfName = 0
    VmNo = 0
    vmRisk = 0
    outputfile = os.path.join(dirpath,"OPCO_OUTPUT")
    for filename in glob.glob(os.path.join(outputfile,opconame +'_'+site+ '_'+ 'AVG'+'__SUMMARY.xlsx')):
        print(filename+ 'File found')
        print('creating presentation slide for '+ filename )
        file_location = filename
        try:            
            multi_sheet_file = pd.ExcelFile(file_location) 
            excel_sheet_names = multi_sheet_file.sheet_names
            for sheetname in excel_sheet_names:
                org_sheetname = sheetname
                print("executing sheet:", sheetname)
                if sheetname == 'Contention_&_Ready_Time':
                    con_df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                    continue
                elif sheetname == 'CPU_Co-Stop':
                    costop_df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                    df = pd.merge(con_df,costop_df,on='Unnamed: 0',how ='inner')
                    df_list = DivideDf(df)
                elif sheetname == 'VMs_at_Risk':
                    df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                    df_list = DivideDf(df)
                    vnfName = df[df.columns[0]]
                    VmNo = df[df.columns[1]]
                    vmRisk = df[df.columns[2]]
                    
                elif sheetname == 'Latency_&_Packet_Drop':
                    df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                    df_list = DivideDf(df)
                    
                else:
                    df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                    df_list = DivideDf(df)
                    
                #limit per slide:15 rows
                currentSlide = 0
                fullsheetname = ''
                for df in df_list:
                    if org_sheetname == 'CPU_Co-Stop':
                        df = pd.DataFrame(np.row_stack([df.columns, df.values]),columns=['1','2','3','4','5','6','7','8','9','10'])
                    elif org_sheetname == 'VMs_at_Risk':
                        df = pd.DataFrame(np.row_stack([df.columns, df.values]),columns=['1','2','3','4','5','6','7','8','9'])
                    elif org_sheetname == 'Latency_&_Packet_Drop':
                        df = pd.DataFrame(np.row_stack([df.columns, df.values]),columns=['1','2','3','4','5','6','7'])
                    else:
                        df = pd.DataFrame(np.row_stack([df.columns, df.values]),columns=['1','2','3','4','5','6'])
                    cols = len(df.columns)
                    rows = len(df.index)
                    #add lside no to title
                    totalslide = len(df_list)
                    currentSlide = currentSlide + 1
                    slideNo = ''
                    if len(df_list)>1:
                        slideNo = ' '+'['+ str(currentSlide)+'/'+str(totalslide)+']'
                    c = cols
                    r = rows
                    if c > 0:
                        slide =  prs.slides.add_slide( prs.slide_layouts[5])
                        shapes = slide.shapes
                        title = slide.shapes.title
                        
                        for i in sheetMap:
                            if sheetname == i:
                                fullsheetname = sheetMap[i] + slideNo
                                
                        title.text = fullsheetname.upper()
                        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(250,0,0)
                        title.text_frame.paragraphs[0].font.size=Pt(32)
                        title.text_frame.paragraphs[0].font.bold=True 
                        left = Inches(0.7)
                        right = Inches(0.1)
                        top = Inches(1.7)
                        width = Inches(4.0)
                        #num = 12.0/c
                        num = 9.6/(c-1)
                        table = shapes.add_table(rows, cols, left, top, width, right).table
                        for i in range(0,c):
                            if i==0:
                                table.columns[i].width = Inches(2.4)
                            else:
                                table.columns[i].width = Inches(num)
                        for i in range(0,r):
                            for e in range(0,c):
                                if i==0 and e==0:
                                    table.cell(i,e).text = 'VNF'
                                else:
                                    table.cell(i,e).text = str(df.iloc[i,e])
                                if e==0 and not i==0:#To adjust the VNF names in 10pt font
                                    cell = table.rows[i].cells[e]
                                    paragraph = cell.text_frame.paragraphs[0]
                                    if len(str(df.iloc[i,e])) > len('ALF-vNeLSALF-vNeLSALF-vNeLSALF-LSA'): 
                                        paragraph.font.size = Pt(8)
                                    else:
                                        paragraph.font.size = Pt(10)
                                    paragraph.alignment = PP_ALIG.CENTER
                                else:
                                    cell = table.rows[i].cells[e]
                                    paragraph = cell.text_frame.paragraphs[0]
                                    paragraph.font.size = Pt(12)
                                    paragraph.alignment = PP_ALIG.CENTER
                                if e==0:
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    fill.fore_color.rgb =  RGBColor(250,204,204)#pink color
                                    
                                if i== 0:
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    fill.fore_color.rgb =  RGBColor(250,0,0)#red color header
                                    
                                if sheetname in ('CPU','Memory'):
                                    if e in (1,4) and not i== 0:
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        fill.fore_color.rgb =  RGBColor(242,242,242)#off-white
                                    if i!=0 and e in (2,3,5):                                    
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        tup =  CPURAMColorScaleRule(df.iloc[i,e])
                                        fill.fore_color.rgb =  getRGBcolor(tup)
                                elif sheetname in ('Latency_&_Packet_Drop'):
                                    if i!=0 and e in (1,2,3):
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        sheetname = 'Latency'
                                        tup =  OtherColorScaleRule(sheetname ,(df.iloc[i,e]))
                                        fill.fore_color.rgb =  getRGBcolor(tup)
                                    elif i!=0 and e in (4,5,6):
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        sheetname = 'Packet_Drop'
                                        tup =  OtherColorScaleRule(sheetname ,(df.iloc[i,e]))
                                        fill.fore_color.rgb =  getRGBcolor(tup)
                                elif sheetname in('Contention_&_Ready_Time','CPU_Co-Stop'):
                                    if i!=0 and e!=0 :
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        tup =  OtherColorScaleRule(sheetname ,(df.iloc[i,e]))
                                        fill.fore_color.rgb =  getRGBcolor(tup)
                                elif sheetname == 'VMs_at_Risk':
                                    if e ==1 and not i== 0:
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        fill.fore_color.rgb =  RGBColor(242,242,242)
                                    if i!=0 and not e in (0,1) :
                                        fill = cell.fill #fill the legend as well
                                        fill.solid()
                                        tup =  OtherColorScaleRule(sheetname ,(df.iloc[i,e]))
                                        fill.fore_color.rgb =  getRGBcolor(tup)
                                
        except Exception as e:
            print("Error!" )
            print(e)
    #CPU at risk Slide
        CPUChartSlide(vnfName,VmNo,vmRisk,prs)
    #heatmap Slide        
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    title = slide.shapes.title
    title.text = "HEATMAPS"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(250,0,0)
    title.text_frame.paragraphs[0].font.size=Pt(32)
    title.text_frame.paragraphs[0].font.bold=True 
    #textbox#
    left =Inches(0.7)
    width = height = Inches(1)
    top = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = 'Here you can find a link to the heatmaps created for each KPI.'+'\n'+'Heatmaps are created for both Avg and Peak hourly data, where peak hourly refers to the peak sample '+'\n'+'in the hour and Avg hourly refers to the Average amongst all 12 ( 5min) samples in the hour.'
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    
    left = Inches(0.7)
    right = Inches(0.5)
    top = Inches(3.3)
    width = Inches(4.0)
    rows=3
    cols = 2
    table = shapes.add_table(rows, cols, left, top, width, right).table 
    
    column_header = [' ','HeatMaps']
    for i in range(0,cols):
        table.columns[i].width = Inches(1.8)
        cell = table.cell(0,i)
        cell.text
        cell.text = column_header[i]
    for i in range(0,cols):
        for e in range(0,rows):
            if e == 0:
                cell = table.rows[e].cells[i]
                fill = cell.fill #fill the legend as well
                fill.solid()
                fill.fore_color.rgb =  RGBColor(250,0,0)
            else:
                cell = table.rows[e].cells[i]
                fill = cell.fill #fill the legend as well
                fill.solid()
                fill.fore_color.rgb =  RGBColor(250,204,204)
    cell = table.cell(1,0) 
    cell.text 
    cell.text = 'AVG'
    cell = table.cell(2,0) 
    cell.text 
    cell.text = 'PEAK'
 
def GenerateCsv(dirpath,opconame,site,maxavg):
    if not os.path.exists(os.path.join(dirpath,'CSV')):
        os.mkdir(os.path.join(dirpath,'CSV'))
    file_list = get_excel_files(OutputFolder)
    for file in file_list:
        if file == opconame +'_'+site+'_'+ maxavg +'__SUMMARY.xlsx':
            multi_sheet_file = pd.ExcelFile(os.path.join(OutputFolder ,file) )   
            excel_sheet_names = multi_sheet_file.sheet_names
            for sheetname in excel_sheet_names:
                df = pd.read_excel(multi_sheet_file, sheet_name=sheetname)
                if sheetname == 'CPU':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'BH',df.columns[2]:'BH_CPU_AVG_VM',df.columns[3]:'BH_CPU_VM',df.columns[4]:'MCBH_CPU_HOUR',df.columns[5]:'MCBH_CPU_VALUE'}, inplace = True)
                elif sheetname == 'Memory':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'BH',df.columns[2]:'BH_MEM_AVG_VM',df.columns[3]:'BH_MEM_VM',df.columns[4]:'MCBH_MEM_HOUR',df.columns[5]:'MCBH_MEM_VALUE'}, inplace = True)
                elif sheetname == 'Contention_&_Ready_Time':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'BH_CPUCON_AVG_VM',df.columns[2]:'BH_CPUCON_VM',df.columns[3]:'MCBH_CPUCON_VALUE',df.columns[4]:'BH_CPURDY_AVG_VM',df.columns[5]:'BH_CPURDY_VM',df.columns[6]:'MCBH_CPURDY_VALUE'}, inplace = True)
                elif sheetname == 'Latency_&_Packet_Drop':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'BH_SRGLTC_AVG_VM',df.columns[2]:'BH_SRGLTC_VM',df.columns[3]:'MCBH_SRGLTC_VALUE',df.columns[4]:'BH_NWPCT_AVG_VM',df.columns[5]:'BH_NWPCT_VM',df.columns[6]:'MCBH_NWPCT_VALUE'}, inplace = True)
                elif sheetname == 'CPU_Co-Stop':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'BH_COSTOP_AVG_VM',df.columns[2]:'BH_COSTOP_VM',df.columns[3]:'MCBH_COSTOP_VALUE'}, inplace = True)
                elif sheetname == 'VMs_at_Risk':
                    df.rename(columns={df.columns[0]:'VMNAME',df.columns[1]:'TOTAL_VM',df.columns[2]:'CPU_ABV_80PCT',df.columns[3]:'MEM_ABV_80PCT',df.columns[4]:'CON_ABV_5PCT',df.columns[5]:'COSTOP_ABV_5PCT',df.columns[6]:'CORDY_ABV_5PCT',df.columns[7]:'SRG_ABV_15MS',df.columns[8]:'NWPCT_ABV_1PCT'}, inplace = True)
                
                df.to_csv(os.path.join(os.path.join(dirpath,'CSV') , opconame +'_'+site+ '_'+ maxavg+'_'+sheetname+'_SUMMARY.csv'),index=False)
opco_site = ''  
def chunksize(n_workers, len_vnlist, factor):#(CPU_Count,len_vnf, 4)
    chunksize, extra = divmod(len_vnlist, n_workers * factor)
    if extra:
        chunksize += 1
    return chunksize
propertyfolder = os.path.join(dirpath,'PropertiesFolder')
settingdata = json.loads(open(os.path.join(propertyfolder,'settings.json')).read())

if bool(settingdata):
    for key,value in settingdata.items():
        if key== 'ExcludeLSFVnf':
            LSF_VM_flag = value
        elif key== 'WeightedVNF':
            WeightedVNF = value
        elif key== 'vTAPs':
            vTAPsFlag = value
def move_slide( presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
# def readGenerator(path):
#     #pd.read_csv(os.path.join(path , file))
#     for row in open(path ):
#         yield row
def main():
    start_time = time.time()    
    path = sys.argv[1]
    map_path = sys.argv[2]
    dsbase_path = sys.argv[3]
    
    logfile = open(os.path.join( OutputFolder,"log"+".txt"),"w")
    Mapping_file = get_excel_files(map_path)
    for map_file in Mapping_file: #if in future we need to take multiple files
        mapping_dataframe = pd.read_csv(os.path.join(map_path,map_file))
        break 
    #mapping_dataframe['VM'] = mapping_dataframe['VM'].fillna('NO_VM') 
    #mapping_dataframe['VNF'] = mapping_dataframe['VNF'].fillna('NO_VNF')
    #mapping_dataframe['Program'] = mapping_dataframe['Program'].fillna('NA') 
    #mapping_dataframe['Vendor'] = mapping_dataframe['Vendor'].fillna('NA')
    #mapping_dataframe['SiteName'] = mapping_dataframe['SiteName'].fillna('NA') 
             
    #mapping_dataframe['VM'] = mapping_dataframe['VM'].map(lambda x:x.split('(')[0])#removing UUID from VM name
    #mapping_dataframe['VM'] = mapping_dataframe['VM'].map(lambda x:x.strip())
    #mapping_dataframe['VNF'] = mapping_dataframe['VNF'].map(lambda x:x.strip())
    #mapping_dataframe['Program'] = mapping_dataframe['Program'].map(lambda x:x.strip())
    #mapping_dataframe['Vendor'] = mapping_dataframe['Vendor'].map(lambda x:x.strip())
    #mapping_dataframe['SiteName'] = mapping_dataframe['SiteName'].map(lambda x:x.strip())
    
    dsbase_file_list = get_excel_files(dsbase_path)
    if not len(dsbase_file_list) == 0:
        for dsbase in dsbase_file_list: #if in future we need to take multiple files
            dsbase_dataframe = pd.read_csv(os.path.join(dsbase_path,dsbase))
            opco_site = dsbase.split('_')[0]
            break
        
        Include_VM_set = set(mapping_dataframe.loc[mapping_dataframe['Include_Manual'].str.lower() == 'yes', 'VM'].tolist())
        #only include YES Vms
        ds_mapping_dataframe = mapping_dataframe[mapping_dataframe['VM'].isin(Include_VM_set)]
        #Keep common VMs which are present in Resname and Sitename.
        ds_mapping_dataframe=ds_mapping_dataframe.rename(columns = {'SiteName':'resName'})
        #ds_mapping_dataframe = ds_mapping_dataframe['resName'].str.lower()
        #dsbase_dataframe = dsbase_dataframe['resName'].str.lower()
        DSBase_DF = pd.merge(ds_mapping_dataframe,dsbase_dataframe, on='resName', how='inner')
        resName_set = set(DSBase_DF['resName'])
        values = []
        value_TC_stg=0.0
        value_P_stg=0.0
        value_O_stg=0.0
            
        value_TC_cpu=0.0
        value_P_cpu=0.0
        value_O_cpu=0.0
        
        value_TC_mem=0.0
        value_P_mem=0.0
        value_O_mem=0.0
        for resName in resName_set:
            temp_dsbase_dataframe = dsbase_dataframe[dsbase_dataframe['resName'] == resName]
            timeStamp = temp_dsbase_dataframe['timestamp'].max()
            temp1_dsbase_dataframe = temp_dsbase_dataframe[temp_dsbase_dataframe['timestamp'] == timeStamp]
            value_TC_stg = value_TC_stg + temp1_dsbase_dataframe.iloc[0,2]#Datastore - Total Capacity (GB)
            value_P_stg = value_P_stg + temp1_dsbase_dataframe.iloc[0,3]
            value_O_stg = value_O_stg + temp1_dsbase_dataframe.iloc[0,4]
            
            value_TC_cpu = value_TC_cpu + temp1_dsbase_dataframe.iloc[0,5]#Datastore - Total Capacity (GB)
            value_P_cpu = value_P_cpu + temp1_dsbase_dataframe.iloc[0,6]
            value_O_cpu = value_O_cpu + temp1_dsbase_dataframe.iloc[0,7]
            
            value_TC_mem = value_TC_mem + temp1_dsbase_dataframe.iloc[0,8]#Datastore - Total Capacity (GB)
            value_P_mem = value_P_mem + temp1_dsbase_dataframe.iloc[0,9]
            value_O_mem = value_O_mem + temp1_dsbase_dataframe.iloc[0,10]
            
        values_stg=[value_TC_stg,value_P_stg,value_O_stg]
        values_cpu=[value_TC_cpu,value_P_cpu,value_O_cpu]
        values_mem=[value_TC_mem,value_P_mem,value_O_mem]
        if value_O_stg == 0.0 or value_O_cpu == 0.0 or value_O_mem == 0.0:
            print("INFO: There is no Vm selected for DSBase bar charts. Please select the VMs as 'YES' in mapping file - Include-Mappng column")
            logfile.write("INFO: There is no Vm selected for DSBase bar charts. Please select the VMs as 'YES' in mapping file - Include-Mappng column")
        else:
            ppt(values_stg,'Storage_Chart',opco_site,'Storage')
            ppt(values_cpu,'CPU_Chart',opco_site,'CPU_Demand')
            ppt(values_mem,'Memory_Chart',opco_site,'Memory')
    else:
        print("INFO: There is no DSBASE file. So bar charts could not be generated."+'\n')
        logfile.write("INFO: There is no DSBASE file. So bar charts could not be generated"+'\n')
    #for DSBase  END
    exclude_VM_set = set(mapping_dataframe.loc[mapping_dataframe['Include_Manual'].str.lower() == 'no', 'VM'].tolist())
    PowerStateOff_set = set(mapping_dataframe.loc[mapping_dataframe['CFG_PowerState'].str.lower() == 'poweredoff', 'VM'].tolist())
    
    LSF_VM_set = {}
    if LSF_VM_flag == True :
        LSF_VM_set = set(mapping_dataframe.loc[mapping_dataframe['CFG_LSF'].str.lower() == 'high', 'VM'].tolist())
        print("INFO: The script will remove all the High Latency VMs")
        logfile.write("INFO: The script will remove all the High Latency VMs"+'\n')
        logfile.write(" "+'\n')
    map_ds_vm = mapping_dataframe[['VM','VNF','SiteName']]

    file_list = get_excel_files(path)
    print(file_list)
    logfile.write("Executing files are "'%s\n' %file_list +'\n')

    site_sum_dicAvg={}
    site_sum_dicMax = {}
    for file in file_list:
        filename = file.split('.')[0]
        opco_site = file.split('_')[0]
        maxavg = re.split('_',filename)[3].upper()
#         #VTAPS
#         if vTAPsFlag == True:
#             if maxavg == 'MAX':
#                 vTaps(os.path.join(path,file),opco_site)
#                 print(" ")
        print("~~~~~~~~~~*** Working on OPCO " + opco_site + "-"+maxavg+" ***~~~~~~~~~~~~~~~")
        print("INFO: number of CPU on this system is "+str(multiprocessing.cpu_count()))
        logfile.write("INFO: number of CPU on this system is "+str(multiprocessing.cpu_count())+'\n')
        # convert the input csv to a dat frame
        try:
            ds_vm = pd.read_csv(os.path.join(path , file))
        except:
            print("ERROR: "+opco_site+" INPUT FILE ARE EMPTY ***")
            print(" __________________________________________________________________________________")
            sys.exit(1)
#         ds_vm = pd.DataFrame()
#         ds_vm = ds_vm.append(readGenerator(os.path.join(path , file)))
#         
#         print(sys.getsizeof(ds_vm))
        ds_vm=ds_vm.rename(columns = {'vCenter VM Name':'VM'})
        ds_vm = pd.merge(map_ds_vm, ds_vm, on='VM', how='inner')
        
        ds_vm['VM'] = ds_vm['VM'].fillna('NO_VM')        
        ds_vm['VNF'] = ds_vm['VNF'].fillna('NO_VNF')
        ds_vm['SiteName'] = ds_vm['SiteName'].fillna('NO_SITE') 
        
        #ds_vm['VM'] = ds_vm['VM'].map(lambda x:x.split('(')[0])#removing UUID from VM name
        ds_vm['VM'] = ds_vm['VM'].map(lambda x:x.strip())
        ds_vm['VNF'] = ds_vm['VNF'].map(lambda x:x.strip())
        ds_vm['SiteName'] = ds_vm['SiteName'].map(lambda x:x.strip())
        
        ds_vm['timestamp'] = pd.to_datetime(ds_vm['timestamp'], format="%Y-%m-%d_%H:%M:%S")
        
        remove_set= set()
        vm_list = ds_vm['VM'].tolist()
        if not len(PowerStateOff_set)==0 :
            for vm in vm_list :
                for item_vm in PowerStateOff_set:
                    if vm == item_vm: 
                        remove_set.add(vm)
        if not len(exclude_VM_set)==0 :
            for vm in vm_list:
                for item_vm in exclude_VM_set:
                    if vm == item_vm: 
                        remove_set.add(vm)
        if not len(LSF_VM_set)==0:
            for vm in vm_list:
                for item_vm in LSF_VM_set:
                    if vm ==item_vm:
                        remove_set.add(vm)
                        
        
        ds_vm = ds_vm[~ds_vm['VM'].isin(remove_set)]
        
        mapping_dataframe = mapping_dataframe[~mapping_dataframe['VM'].isin(remove_set)]
        vendor_df = mapping_dataframe [['VNF','Program','Vendor']]
        
        datetime_set = set(ds_vm['timestamp'])
        datetime_set = sorted(datetime_set)
        temp_date = datetime.strptime(min(datetime_set).strftime("%Y-%m-%d_00:00:00"), "%Y-%m-%d_%H:%M:%S")
        min_date = temp_date
        max_date = max(datetime_set)
        ds_vm.sort_values(by=['VNF', 'timestamp'])
        month = temp_date.strftime("%B")
        
        set_site = set(ds_vm['SiteName'].str.strip())#get the sites name without duplicate
        sites={}   
        for eachsite in set_site:
            ds_vm_temp = ds_vm[ds_vm['SiteName'] == eachsite]
            vnf_set = set(ds_vm_temp['VNF'])
            print(vnf_set)
            logfile.write("Number of VNFs executing are : " '%s\n' %vnf_set+'\n')
            sites[eachsite] = vnf_set

        print("~~~~~~~~~~~~~~ Number of VNF in sheet "+maxavg+":"+" ~~~~~~~~~~~~~~"+'\n')
        vnfname = set()
        for key, values in sites.items():
            vnfname= vnfname.union(values)
        print(" ")
        logfile.write("~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~" +'\n')
        fullopconame =''
        #slide 1: opoco slide
        prs = Presentation(os.path.join(propertyfolder,'sample.pptx'))
        OPCOMap = {'RO': 'ROMANIA','CZ': 'CZECH REPUBLIC','HU': 'HUNGARY','PT': 'PORTUGAL','ES': 'SPAIN','UK': 'UNITED KINGDOM','IT': 'ITALY','AL': 'ALBANIA','DE': 'GERMANY','GR': 'GREECE','Central':'CENTRAL','IE':'IRELAND'}
        for i in OPCOMap:
            if opco_site.upper() in i:
                fullopconame = OPCOMap[i]
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        title = slide.shapes.title
        title.text = fullopconame
        subtitle=slide.placeholders[1]
        subtitle.text = ' '
        #title.text_frame.fit_text(max_size=32,bold=True)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)  
        title.text_frame.paragraphs[0].font.size=Pt(44)
        
        imagepath = os.path.join(os.path.join(dirpath,"PropertiesFolder"),'Flags')
        if opco_site.upper() in ('RO','CZ','HU','PT','ES','CENTRAL','UK','IT','AL','DE','GR','IE'):
            img_path = os.path.join(imagepath,opco_site+'.png')
        else:
            img_path = os.path.join(imagepath,'ERROR.png')
        pic = slide.shapes.add_picture(img_path, Inches(10.0), Inches(0.5),width=Inches(3), height=Inches(2))
        #feedback slide ~~~~~~~~~~~~~~~~~~~~~~~~~
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        title = slide.shapes.title
        title.text = "FEEDBACK"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(250,0,0)
        title.text_frame.paragraphs[0].font.size=Pt(32)
        title.text_frame.paragraphs[0].font.bold=True 
#         left = Inches(0.5)
#         right = Inches(0.1)
#         top = Inches(1.7)
#         width = Inches(4.0)
#         rows=13
#         cols = 6
#         table = shapes.add_table(rows, cols, left, top, width, right).table 
#         feedback_header = ['VNF','VM# Impacted','Issue','Repeat Issue (Y/N)', 'Design Comment','Action Point']
#         for i in range(0,cols):
#             if i in (0,2,4,5):
#                 table.columns[i].width = Inches(2.7)
#             elif i in (1,3):
#                 table.columns[i].width = Inches(0.7)
#             cell = table.cell(0,i)
#             cell.text = feedback_header[i]
#     
#         for i in range(0,cols):
#             for e in range(0,rows):
#                 if e == 0:
#                     cell = table.rows[e].cells[i]
#                     fill = cell.fill #fill the legend as well
#                     fill.solid()
#                     fill.fore_color.rgb =  RGBColor(250,0,0)#red color
#                     paragraph = cell.text_frame.paragraphs[0]
#                     paragraph.font.size = Pt(10)
#                     paragraph.alignment = PP_ALIG.CENTER
#                 elif i in (0,1) and e!=0 :
#                     cell = table.rows[e].cells[i]
#                     fill = cell.fill #fill the legend as well
#                     fill.solid()
#                     fill.fore_color.rgb =  RGBColor(250,204,204) #pink color
#                 else:
#                     cell = table.rows[e].cells[i]
#                     fill = cell.fill #fill the legend as well
#                     fill.solid()
#                     fill.fore_color.rgb =  RGBColor(222,235,247) #blue color 
        #end of intro slides creating             
        for site, vnf_set in sites.items():
            print("~~~~~~~~~~*** Working on site " + site + " ***~~~~~~~~~~~~~~~"+'\n')
            logfile.write("~~~~~~~~~~*** Working on site " + site + " ***~~~~~~~~~~~~~~~"+'\n')
            Avg_Df =pd.DataFrame()
            Max_Df =pd.DataFrame()
            vnf_num = len(vnf_set)
            
            if (multiprocessing.cpu_count() <=8):#as the office systems will not be able to perform any other tasks.
                cpucount= multiprocessing.cpu_count()-1
            else:
                cpucount = multiprocessing.cpu_count() 
            
            pool = multiprocessing.Pool(cpucount)
            
            result_list = pool.map(partial(executeVNF,maxavg=maxavg,ds_vm=ds_vm,min_date=min_date,max_date=max_date,site=site),vnf_set,chunksize(cpucount, vnf_num, 4)) #
            pool.close()
            pool.join()
            List_of_parameters = []
            for a,b in result_list:
                #heatmap_data = c
                List_of_parameters.append(b)
                for key,val in a.items():
                    if key == 'AVG':
                        Avg_Df = Avg_Df.append(val)
                    elif key =='MAX':
                        Max_Df = Max_Df.append(val)  
            if not Avg_Df.empty:
                site_sum_dicAvg[site] = Avg_Df
            if not Max_Df.empty:
                site_sum_dicMax[site] = Max_Df
            #write to heat map files   
            WriteVnfData(List_of_parameters,opco_site,site,maxavg)
            GenerateCsv(dirpath,opco_site,site,maxavg)
            #Create the presentation slides.
            autoSlide(dirpath, opco_site, site,prs) 
        #Refence slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        title = slide.shapes.title
        title.text = 'Reference'
        subtitle=slide.placeholders[1]
        subtitle.text = ' '
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
        title.text_frame.paragraphs[0].font.size=Pt(32)
        title.text_frame.paragraphs[0].font.bold=True 
        for s in range(0,7):
            move_slide( prs, 2, len(prs.slides))
        
        prs.save(os.path.join(os.path.join(OutputFolder,'VF-'+opco_site.upper()+'-'+'CP Report-'+str(month)+'-'+str(datetime.today().year)+'.pptx')))
    

    MAX_exe_df = pd.DataFrame()
    AVG_exe_df = pd.DataFrame()
    
    list_of_all_datas_avg = pd.DataFrame()
    list_of_all_datas_max = pd.DataFrame()
    list_of_all_datas_avgP = pd.DataFrame()
    list_of_all_datas_maxP= pd.DataFrame()
        
    if bool(site_sum_dicMax) and bool(site_sum_dicAvg):
        for site, Avg_Df1 in site_sum_dicAvg.items():
            Max_Df1 = site_sum_dicMax[site]
            result = executiveSum(vendor_df,Max_Df1,Avg_Df1,opco_site,site)
            logfile.write("per site executive summary"+'\n')
            logfile.write("for "+ site+ "AVG"+'\n')
            logfile.write(str(result[0])+'\n')
            logfile.write("for "+ site+ "MAX"+'\n')
            logfile.write(str(result[2])+'\n')
            
            list_of_all_datas_avg = list_of_all_datas_avg.append(result[0]) 
            list_of_all_datas_avgP = list_of_all_datas_avgP.append(result[1])
            list_of_all_datas_max = list_of_all_datas_max.append(result[2])
            list_of_all_datas_maxP = list_of_all_datas_maxP.append(result[3])
        
        list_of_all_datas_avg = list_of_all_datas_avg.reset_index(drop=True) 
        list_of_all_datas_avgP = list_of_all_datas_avgP.reset_index(drop=True) 
        list_of_all_datas_max = list_of_all_datas_max.reset_index(drop=True) 
        list_of_all_datas_maxP = list_of_all_datas_maxP.reset_index(drop=True) 
                     
    else:
        print("INFO:  Executive summary cannot be generated. Need both MAX and AVG"+'\n')
        logfile.write("INFO:  Executive summary cannot be generated. Need both MAX and AVG"+'\n')
    #aggregaet the executive sum.
    #if not MAX_exe_df.empty and not AVG_exe_df.empty:
    if not list_of_all_datas_avg.empty and not list_of_all_datas_avgP.empty and not list_of_all_datas_max.empty and not list_of_all_datas_maxP.empty:
        #vendor df
        list_of_all_datas_avg = list_of_all_datas_avg.T
        list_of_all_datas_max = list_of_all_datas_max.T
        list_of_all_datas_avg['VALUE'] = list_of_all_datas_avg.mean(numeric_only=True, axis=1)
        list_of_all_datas_max['VALUE'] = list_of_all_datas_max.mean(numeric_only=True, axis=1)
        list_of_all_datas_avg['PGM_VDR'] = list_of_all_datas_avg.index
        list_of_all_datas_max['PGM_VDR'] = list_of_all_datas_max.index
        list_of_all_datas_avg = list_of_all_datas_avg[['PGM_VDR','VALUE']]
        list_of_all_datas_max = list_of_all_datas_max[['PGM_VDR','VALUE']]
        print(list_of_all_datas_avg['VALUE'])
        new_rowA = {'PGM_VDR': 'Overall', 'VALUE': list_of_all_datas_avg['VALUE'].mean( skipna = True)}
        new_rowM = {'PGM_VDR': 'Overall', 'VALUE': list_of_all_datas_max['VALUE'].mean( skipna = True)}
        list_of_all_datas_avg = list_of_all_datas_avg.append(new_rowA, ignore_index=True)
        list_of_all_datas_max = list_of_all_datas_max.append(new_rowM, ignore_index=True)
        
        #program df
        list_of_all_datas_avgP = list_of_all_datas_avgP.T
        list_of_all_datas_maxP = list_of_all_datas_maxP.T
        list_of_all_datas_avgP['VALUE'] = list_of_all_datas_avgP.mean(numeric_only=True, axis=1)
        list_of_all_datas_maxP['VALUE'] = list_of_all_datas_maxP.mean(numeric_only=True, axis=1)
        list_of_all_datas_avgP['PGM_VDR'] = list_of_all_datas_avgP.index
        list_of_all_datas_maxP['PGM_VDR'] = list_of_all_datas_maxP.index
        list_of_all_datas_avgP = list_of_all_datas_avgP[['PGM_VDR','VALUE']]
        list_of_all_datas_maxP = list_of_all_datas_maxP[['PGM_VDR','VALUE']]

        AVG_exe_df = list_of_all_datas_avg.append(list_of_all_datas_avgP)
        MAX_exe_df = list_of_all_datas_max.append(list_of_all_datas_maxP)
        AVG_exe_df['OPCO'] = opco_site
        MAX_exe_df['OPCO'] = opco_site
                
        #generate excel for executive sum.
        writer = ExcelWriter(os.path.join(OutputFolder, opco_site + '_'+'ExecutiveSummary.xlsx'))
        MAX_exe_df.to_excel(writer, sheet_name='MAX',index=False) #MAX
        AVG_exe_df.to_excel(writer, sheet_name='AVG',index=False)#AVG
        writer.save()
        #csv
        MAX_exe_df.to_csv(os.path.join(os.path.join(dirpath,'CSV') , opco_site +'_'+'MAX'+ '_ExecutiveSummary.csv'),index=False) #MAX
        AVG_exe_df.to_csv(os.path.join(os.path.join(dirpath,'CSV') , opco_site +'_'+'AVG'+ '_ExecutiveSummary.csv'),index=False)#AVG
    else:
        print("INFO:  Executive summary cannot be generated due to some issue with data")
        logfile.write("INFO:  Executive summary cannot be generated due to some issue with data"+'\n')
    end_tim = time.time()
    print('Execution Time:' + '%.2f' % ((end_tim - start_time)/60) +' minutes')
    logfile.write('Execution Time:' + '%.2f' % ((end_tim - start_time)/60) +' minutes'+'\n')
    logfile.close()
if __name__ == '__main__':
    main()
